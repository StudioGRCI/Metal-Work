# -*- coding: utf-8 -*-
"""
Arma la presentación del procedimiento de la orden de trabajo con las capturas
reales del sistema. Se ejecuta después de capturar-procedimiento.mjs.
"""
import os
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

CAPTURAS = os.environ.get('CAPTURAS', '/tmp/presentacion')
RECORTES = os.environ.get('RECORTES', '/tmp/slides')
SALIDA = os.environ.get('SALIDA_PPTX', '/tmp/Procedimiento-Orden-de-Trabajo.pptx')
RAIZ = os.path.abspath(os.path.join(os.path.dirname(__file__), '..', '..'))

AZUL     = RGBColor(0x10, 0x47, 0x80)
AZUL_OSC = RGBColor(0x0A, 0x2B, 0x4F)
ROJO     = RGBColor(0xE5, 0x1D, 0x20)
BLANCO   = RGBColor(0xFF, 0xFF, 0xFF)
CLARO    = RGBColor(0xC5, 0xD6, 0xE8)
GRIS     = RGBColor(0x5B, 0x6B, 0x7B)
TINTA    = RGBColor(0x12, 0x1C, 0x28)

ANCHO, ALTO = Inches(13.333), Inches(7.5)
BANDA = Inches(1.02)

# etapa, archivo, título, explicación
PASOS = [
    ('PUNTO DE PARTIDA', '01-tablero', 'El tablero del taller',
     'Lo que está abierto, en proceso, pausado, atrasado y urgente. Es la pantalla con la que se empieza el día.'),

    ('1 · SE COTIZA', '02-cotizacion', 'La cotización',
     'Cliente, unidad, plazo y precio. De acá salen el presupuesto de la orden y la ficha del producto.'),
    ('1 · SE COTIZA', '03-ficha-tecnica', 'La ficha técnica del producto',
     'Espesores, normas de soldadura y sistemas, como los escribe la empresa. Se trae de una ficha ya redactada y solo se ajusta lo que cambia.'),
    ('1 · SE COTIZA', '04-accesorios-cot', 'Los accesorios que se ofrecen',
     'Lo que la cotización promete entregar. El que trae solo el soporte queda marcado como tal, y así viaja hasta la entrega.'),

    ('2 · SE ABRE LA ORDEN', '05-ordenes', 'Las órdenes del taller',
     'Todas las órdenes con su estado, su avance y su fecha comprometida.'),
    ('2 · SE ABRE LA ORDEN', '06-orden-nueva', 'Se abre la orden',
     'Desde la cotización aprobada, con el cliente, la unidad y el presupuesto ya arrastrados. La orden nace en borrador.'),
    ('2 · SE ABRE LA ORDEN', '07-orden-resumen', 'La orden de trabajo',
     'Avance, horas, presupuesto y entrega comprometida en una sola línea. Debajo, todo lo demás por pestañas.'),

    ('3 · SE APRUEBA', '08-etapas', 'Las etapas de fabricación',
     'Al aprobar la orden se instancian las etapas del catálogo con sus horas estándar. El avance de la orden se pondera con ellas.'),
    ('3 · SE APRUEBA', '09-ficha-medidas', 'Ficha de taller · las medidas',
     'Sección 4 del formato: largo, ancho, alto, capacidad, ruedas, ejes, suspensión y colores. Lo que el taller necesita a la vista.'),
    ('3 · SE APRUEBA', '10-verificacion', 'Ficha de taller · la verificación',
     'Sección 11: los pasos que hay que recorrer antes de dar la unidad por terminada, cada uno con su responsable y sus dos pasadas.'),
    ('3 · SE APRUEBA', '11-accesorios-ot', 'Ficha de taller · los accesorios',
     'Sección 6: los accesorios bajan de la cotización. Lo que se prometió es lo que hay que montar, y cada uno se marca al montarse.'),

    ('4 · SE FABRICA', '12-requerimiento', 'El requerimiento de material',
     'El taller pide lo que necesita para la orden. Almacén lo atiende con lo que hay o lo manda a comprar.'),
    ('4 · SE FABRICA', '13-movimiento', 'El movimiento de almacén',
     'Cada salida queda cargada a la orden y valorizada al costo del material. Es lo que después aparece en el costo real.'),
    ('4 · SE FABRICA', '14-parte-diario', 'El parte diario de producción',
     'Horas por operario, orden y etapa. Al aprobar el parte, las horas se cargan a la orden con el costo hora de cada uno.'),
    ('4 · SE FABRICA', '15-avance-taller', 'El avance en taller',
     'Una tarjeta por unidad: en qué etapa está, hace cuánto no se toca y qué la traba.'),
    ('4 · SE FABRICA', '16-avance-unidad', 'El avance del día, con fotos',
     'Lo que se hizo hoy en la unidad, con las fotos de la jornada. Registrar un avance mueve la etapa.'),
    ('4 · SE FABRICA', '17-servicios', 'El trabajo que se manda afuera',
     'Órdenes de servicio al subcontratista, con su plazo, su conformidad y su pago. También carga a la orden.'),

    ('5 · SE VERIFICA', '18-calidad', 'Control de calidad',
     'Inspecciones con su resultado, sus observaciones y el levantamiento de cada una.'),
    ('5 · SE VERIFICA', '19-documentos', 'Los documentos de la orden',
     'Planos, actas y certificados, versionados y con descarga por enlace temporal.'),
    ('5 · SE VERIFICA', '20-firmas', 'Las firmas pendientes',
     'Cada quien ve lo que espera su firma. El documento no queda aprobado hasta que la cadena se completa.'),

    ('6 · SE ENTREGA', '21-costos', 'Costo real contra presupuesto',
     'Material, mano de obra y servicios de la orden, contra lo que se presupuestó al cotizar.'),
    ('6 · SE ENTREGA', '22-trazabilidad', 'La trazabilidad',
     'Todo lo que le pasó a la orden, quién lo hizo y cuándo. No se puede borrar.'),
    ('6 · SE ENTREGA', '23-informes', 'Los informes',
     'Producción, cumplimiento de entregas, rentabilidad, cotizaciones, consumo de material y subcontratos.'),
    ('6 · SE ENTREGA', '24-cumplimiento', 'Cumplimiento de entregas',
     'Cuántas unidades salieron a tiempo y cuántos días se pasaron las que no. El período se elige arriba.'),
]

RECORRIDO = [
    ('1', 'Se cotiza',       'Ficha técnica del producto, accesorios ofrecidos y precio.'),
    ('2', 'Se abre la orden','Desde la cotización aprobada, con el presupuesto arrastrado.'),
    ('3', 'Se aprueba',      'Se crean las etapas y se arma la ficha de taller: medidas, accesorios y verificación.'),
    ('4', 'Se fabrica',      'Material, horas, avance con fotos y trabajo mandado afuera.'),
    ('5', 'Se verifica',     'Calidad, documentos firmados y los pasos de verificación al 100%.'),
    ('6', 'Se entrega',      'Costo real contra presupuesto, trazabilidad e informes.'),
]

os.makedirs(RECORTES, exist_ok=True)
pres = Presentation()
pres.slide_width, pres.slide_height = ANCHO, ALTO
VACIA = pres.slide_layouts[6]


def fondo(diapo, color):
    fondo = diapo.background.fill
    fondo.solid()
    fondo.fore_color.rgb = color


def caja(diapo, x, y, cx, cy, texto, tamano, color, negrita=False,
         alineacion=PP_ALIGN.LEFT, espaciado=None, fuente='Calibri'):
    cuadro = diapo.shapes.add_textbox(x, y, cx, cy)
    marco = cuadro.text_frame
    marco.word_wrap = True
    marco.vertical_anchor = MSO_ANCHOR.TOP
    marco.margin_left = marco.margin_right = marco.margin_top = marco.margin_bottom = 0
    p = marco.paragraphs[0]
    p.alignment = alineacion
    t = p.add_run()
    t.text = texto
    t.font.size = Pt(tamano)
    t.font.bold = negrita
    t.font.color.rgb = color
    t.font.name = fuente
    if espaciado:
        p.line_spacing = espaciado
    return cuadro


def rectangulo(diapo, x, y, cx, cy, color):
    from pptx.enum.shapes import MSO_SHAPE
    f = diapo.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, cx, cy)
    f.fill.solid()
    f.fill.fore_color.rgb = color
    f.line.fill.background()
    f.shadow.inherit = False
    return f


def recortar(nombre, proporcion):
    """Recorta la captura por abajo para que entre completa en la diapositiva."""
    origen = os.path.join(CAPTURAS, nombre + '.png')
    im = Image.open(origen)
    alto = int(im.width / proporcion)
    if alto < im.height:
        im = im.crop((0, 0, im.width, alto))
    destino = os.path.join(RECORTES, nombre + '.png')
    im.save(destino)
    return destino


# ------------------------------------------------------------------- portada
d = pres.slides.add_slide(VACIA)
fondo(d, AZUL_OSC)
rectangulo(d, 0, ALTO - Inches(0.28), ANCHO, Inches(0.28), ROJO)

logo = os.path.join(RAIZ, 'public/marca/logo-metal-work-claro.png')
if os.path.exists(logo):
    d.shapes.add_picture(logo, Inches(1.1), Inches(1.0), height=Inches(1.0))

caja(d, Inches(1.1), Inches(2.6), Inches(11), Inches(1.6),
     'Procedimiento de la\norden de trabajo', 48, BLANCO, True, espaciado=1.05)
caja(d, Inches(1.1), Inches(4.5), Inches(10), Inches(1.0),
     'Cómo se recorre en el sistema, de la cotización a la entrega de la unidad',
     20, CLARO)
caja(d, Inches(1.1), Inches(5.35), Inches(10), Inches(0.6),
     'METAL WORK PERÚ S.A.C.  ·  Capturas tomadas del sistema en funcionamiento',
     13, RGBColor(0x8F, 0xA8, 0xC2))

# ------------------------------------------------------- el recorrido entero
d = pres.slides.add_slide(VACIA)
fondo(d, BLANCO)
rectangulo(d, 0, 0, ANCHO, BANDA, AZUL)
caja(d, Inches(0.75), Inches(0.28), Inches(11), Inches(0.6),
     'El recorrido, en una página', 26, BLANCO, True)

y = Inches(1.55)
for numero, titulo, detalle in RECORRIDO:
    rectangulo(d, Inches(0.75), y, Inches(0.52), Inches(0.52), AZUL)
    caja(d, Inches(0.75), y + Inches(0.10), Inches(0.52), Inches(0.4),
         numero, 16, BLANCO, True, PP_ALIGN.CENTER)
    caja(d, Inches(1.55), y + Inches(0.01), Inches(3.1), Inches(0.4),
         titulo, 17, TINTA, True)
    caja(d, Inches(4.8), y + Inches(0.05), Inches(7.8), Inches(0.5),
         detalle, 13, GRIS)
    y += Inches(0.93)

rectangulo(d, Inches(0.75), Inches(7.0), Inches(11.83), Inches(0.03),
           RGBColor(0xE2, 0xE8, 0xEF))

# ---------------------------------------------------------- una por pantalla
imagen_ancho = ANCHO - Inches(1.1)
imagen_alto = ALTO - BANDA - Inches(0.55)
proporcion = Emu(imagen_ancho).inches / Emu(imagen_alto).inches

for etapa, archivo, titulo, explicacion in PASOS:
    if not os.path.exists(os.path.join(CAPTURAS, archivo + '.png')):
        print('falta la captura', archivo)
        continue

    d = pres.slides.add_slide(VACIA)
    fondo(d, RGBColor(0xF4, 0xF7, 0xFA))
    rectangulo(d, 0, 0, ANCHO, BANDA, AZUL)
    rectangulo(d, 0, BANDA, ANCHO, Inches(0.04), ROJO)

    caja(d, Inches(0.55), Inches(0.13), Inches(4), Inches(0.3),
         etapa, 10.5, RGBColor(0x9F, 0xC0, 0xE4), True)
    caja(d, Inches(0.55), Inches(0.40), Inches(7.4), Inches(0.5),
         titulo, 21, BLANCO, True)
    caja(d, Inches(8.2), Inches(0.24), Inches(4.6), Inches(0.75),
         explicacion, 11, RGBColor(0xD3, 0xE2, 0xF2), espaciado=1.18)

    d.shapes.add_picture(recortar(archivo, proporcion),
                         Inches(0.55), BANDA + Inches(0.30),
                         width=imagen_ancho)

# ------------------------------------------------------------------- cierre
d = pres.slides.add_slide(VACIA)
fondo(d, AZUL_OSC)
rectangulo(d, 0, ALTO - Inches(0.28), ANCHO, Inches(0.28), ROJO)
caja(d, Inches(1.1), Inches(1.2), Inches(11), Inches(0.9),
     'Qué queda verificado al entregar', 34, BLANCO, True)

CIERRE = [
    'Los accesorios que se cotizaron están montados, cada uno con el V°B° de quién lo revisó y cuándo.',
    'Los pasos de verificación de la carrocería están recorridos, con su primera pasada y su revisión.',
    'Las horas y el material cargados a la orden se pueden contrastar contra el presupuesto de la cotización.',
    'Los documentos de la unidad están firmados por quien correspondía, en el orden que correspondía.',
    'Todo lo anterior quedó registrado con fecha y responsable, y no se puede borrar.',
]
y = Inches(2.4)
for linea in CIERRE:
    rectangulo(d, Inches(1.1), y + Inches(0.12), Inches(0.12), Inches(0.12), ROJO)
    caja(d, Inches(1.55), y, Inches(10.6), Inches(0.7), linea, 15, CLARO, espaciado=1.2)
    y += Inches(0.8)

pres.save(SALIDA)
print('✔', SALIDA, '·', len(pres.slides.__iter__.__self__._sldIdLst), 'diapositivas')


# =============================================================================
# La misma presentación en HTML, para imprimirla en PDF.
# En esta máquina LibreOffice viene sin Impress, así que el PDF no puede salir
# del .pptx: se arma aparte con el mismo contenido y el mismo diseño.
# =============================================================================
import base64
import html as _html

SALIDA_HTML = os.environ.get('SALIDA_HTML', '/tmp/procedimiento.html')


def incrustar(ruta):
    with open(ruta, 'rb') as f:
        return 'data:image/png;base64,' + base64.b64encode(f.read()).decode()


e = _html.escape
partes = []

logo_datos = incrustar(logo) if os.path.exists(logo) else None

partes.append(f"""<section class="lamina portada">
  {'<img class="logo" src="' + logo_datos + '" alt="Metal Work">' if logo_datos else ''}
  <h1>Procedimiento de la<br>orden de trabajo</h1>
  <p class="bajada">Cómo se recorre en el sistema, de la cotización a la entrega de la unidad</p>
  <p class="pie">METAL WORK PERÚ S.A.C. · Capturas tomadas del sistema en funcionamiento</p>
</section>""")

filas = ''.join(
    f'<li><span class="numero">{n}</span><span class="titulo">{e(t)}</span>'
    f'<span class="detalle">{e(d)}</span></li>'
    for n, t, d in RECORRIDO)
partes.append(f"""<section class="lamina">
  <header class="banda"><h2>El recorrido, en una página</h2></header>
  <ol class="recorrido">{filas}</ol>
</section>""")

for etapa, archivo, titulo, explicacion in PASOS:
    recorte = os.path.join(RECORTES, archivo + '.png')
    if not os.path.exists(recorte):
        continue
    partes.append(f"""<section class="lamina">
  <header class="banda banda-paso">
    <div class="izquierda">
      <p class="etapa">{e(etapa)}</p>
      <h2>{e(titulo)}</h2>
    </div>
    <p class="explicacion">{e(explicacion)}</p>
  </header>
  <div class="captura"><img src="{incrustar(recorte)}" alt="{e(titulo)}"></div>
</section>""")

cierre = ''.join(f'<li>{e(linea)}</li>' for linea in CIERRE)
partes.append(f"""<section class="lamina portada cierre">
  <h1 class="chico">Qué queda verificado al entregar</h1>
  <ul class="lista">{cierre}</ul>
</section>""")

ESTILO = """
@page { size: 13.333in 7.5in; margin: 0 }
* { box-sizing: border-box; margin: 0; padding: 0 }
body { font-family: Calibri, 'Segoe UI', system-ui, sans-serif; -webkit-print-color-adjust: exact;
       print-color-adjust: exact }
.lamina { width: 13.333in; height: 7.5in; page-break-after: always; position: relative;
          overflow: hidden; background: #F4F7FA }
.lamina:last-child { page-break-after: auto }

.portada { background: #0A2B4F; color: #fff; padding: 1.1in }
.portada::after { content: ''; position: absolute; left: 0; right: 0; bottom: 0;
                  height: 0.28in; background: #E51D20 }
.portada .logo { height: 1in; margin-bottom: 0.55in }
.portada h1 { font-size: 48pt; font-weight: 700; line-height: 1.05; letter-spacing: -0.01em }
.portada h1.chico { font-size: 34pt; margin-top: 0.2in }
.portada .bajada { margin-top: 0.5in; font-size: 20pt; color: #C5D6E8 }
.portada .pie { margin-top: 0.35in; font-size: 13pt; color: #8FA8C2 }
.cierre .lista { margin-top: 0.6in; list-style: none }
.cierre .lista li { position: relative; padding-left: 0.45in; margin-bottom: 0.34in;
                    font-size: 15pt; color: #C5D6E8; line-height: 1.35 }
.cierre .lista li::before { content: ''; position: absolute; left: 0; top: 0.12in;
                            width: 0.12in; height: 0.12in; background: #E51D20 }

.banda { background: #104780; color: #fff; height: 1.02in; padding: 0.13in 0.55in;
         border-bottom: 0.04in solid #E51D20 }
.banda h2 { font-size: 26pt; font-weight: 700; margin-top: 0.15in }
.banda-paso { display: flex; align-items: flex-start; justify-content: space-between; gap: 0.5in }
.banda-paso .izquierda { max-width: 7.4in }
.banda-paso .etapa { font-size: 10.5pt; font-weight: 700; color: #9FC0E4; letter-spacing: 0.06em }
.banda-paso h2 { font-size: 21pt; margin-top: 0.05in }
.banda-paso .explicacion { width: 4.6in; font-size: 11pt; color: #D3E2F2; line-height: 1.18;
                           margin-top: 0.11in }

.captura { padding: 0.30in 0.55in 0.25in }
.captura img { width: 100%; display: block }

.recorrido { list-style: none; margin: 0.53in 0.75in }
.recorrido li { display: flex; align-items: baseline; gap: 0.28in; margin-bottom: 0.41in }
.recorrido .numero { flex: 0 0 0.52in; height: 0.52in; background: #104780; color: #fff;
                     font-size: 16pt; font-weight: 700; display: inline-flex;
                     align-items: center; justify-content: center }
.recorrido .titulo { flex: 0 0 3.1in; font-size: 17pt; font-weight: 700; color: #121C28 }
.recorrido .detalle { flex: 1; font-size: 13pt; color: #5B6B7B }
"""

with open(SALIDA_HTML, 'w') as f:
    f.write('<!doctype html><html lang="es"><head><meta charset="utf-8">'
            '<title>Procedimiento de la orden de trabajo</title>'
            f'<style>{ESTILO}</style></head><body>' + ''.join(partes) + '</body></html>')

print('✔', SALIDA_HTML)
